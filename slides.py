"""
Adapted from makeslides.py from the "User Stories Template Tool" created by a CSSE3002
student in 2018.
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


prs = Presentation()

SLD_LAYOUT_BLANK = 6
BULLET  = chr(8226)


def create_slides(stories, filename):
    for story in stories:

        slide_layout = prs.slide_layouts[SLD_LAYOUT_BLANK]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        left = Cm(0.21)
        top = Cm(0.24)
        width = Cm(2.49)
        height = Cm(1.93)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        fill.fore_color.brightness = 0.4

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = 'Story ID: ' + story.id_

        ######################### STORY_TITLE
        left = Cm(3.02)
        top = Cm(0.24)
        width = Cm(15.62)
        height = Cm(2.05)

        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1


        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)
        text_frame.word_wrap = True


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = story.title

        ######################### PRIORITY

        left = Cm(18.96)
        top = Cm(0.24)
        width = Cm(3.01)
        height = Cm(1.93)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        fill.fore_color.brightness = 0.4

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = 'Priority: '+ story.priority

        ######################### STORY POINTS


        left = Cm(22.1)
        top = Cm(0.24)
        width = Cm(3.16)
        height = Cm(1.93)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        fill.fore_color.brightness = 0.4

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = 'Story Points: ' + story.points

        ######################### STORY


        left = Cm(0.21)
        top = Cm(2.41)
        width = Cm(24.98)
        height = Cm(7.4)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(109,158,235)

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)
        text_frame.margin_left = Cm(0.13)
        text_frame.margin_right = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.LEFT
        p.text = '\t'+ story.body


        ######################### ACCEPTANCE CRITERIA


        left = Cm(0.21)
        top = Cm(10.03)
        width = Cm(24.98)
        height = Cm(4.67)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(109,158,235)

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)
        text_frame.margin_left = Cm(0.13)
        text_frame.margin_right = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = 'Acceptance Criteria:'


        p = text_frame.add_paragraph()
        p.font.color.rgb = RGBColor(0,0,0)

        paragraph_strs = story.criteria.split('\n')
        p.text = '\t'+BULLET+' '+paragraph_strs[0]
        for para_str in paragraph_strs[1:]:
            p = text_frame.add_paragraph()
            p.font.color.rgb = RGBColor(0,0,0)
            p.text = '\t'+BULLET+' '+para_str

        ######################### NOTES


        left = Cm(0.21)
        top = Cm(14.92)
        width = Cm(24.98)
        height = Cm(3.64)


        shape = shapes.add_shape(
           MSO_SHAPE.RECTANGLE, left, top, width, height
           )

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,255,255)

        line = shape.line
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        line.color.brightness = -0.5 
        line.width = Pt(.75)

        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)
        text_frame.margin_left = Cm(0.13)
        text_frame.margin_right = Cm(0.13)


        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.bold = False
        p.font.color.rgb = RGBColor(0,0,0)
        p.text = 'Notes:'
        p = text_frame.add_paragraph()
        p.font.color.rgb = RGBColor(0,0,0)
        

        paragraph_strs = story.notes.split('\n')
        p.text = '\t'+BULLET+' '+paragraph_strs[0]
        for para_str in paragraph_strs[1:]:
            p = text_frame.add_paragraph()
            p.font.color.rgb = RGBColor(0,0,0)
            p.text = '\t'+BULLET+' '+para_str


    prs.save(filename)